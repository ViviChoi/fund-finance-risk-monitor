"""
Generate slides/FFRM-Deck.pptx -- a 12-slide banker-friendly walkthrough
of the Fund Finance Facility Risk Monitor.

Run:
    python build_deck.py

Requirements: python-pptx (in requirements.txt).
"""

import os
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FIGURES = ROOT / "figures"
SLIDES = ROOT / "slides"
sys.path.insert(0, str(ROOT))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
GOLD = RGBColor(0xE5, 0xA8, 0x39)
GREY_DARK = RGBColor(0x2F, 0x3B, 0x4A)
GREY_LIGHT = RGBColor(0x80, 0x8B, 0x9C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xE6, 0x91, 0x38)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def add_blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def add_text_box(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=GREY_DARK, align=PP_ALIGN.LEFT,
                 font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=GREY_DARK, line_spacing=1.2,
                bullet_color=NAVY):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_left = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "▸  " + item
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tx


def add_header_band(slide, title, subtitle=None):
    # navy band on top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
                 title, size=26, bold=True, color=WHITE, font=FONT_HEAD)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                     subtitle, size=14, color=GREY_LIGHT)


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        print(f"  ! missing image: {path}")
        return None
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width, height)
    elif width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    elif height:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    else:
        return slide.shapes.add_picture(str(path), left, top)


def add_table(slide, left, top, width, height, data, *,
              header_fill=NAVY, header_color=WHITE,
              row_fill=None, alt_fill=RGBColor(0xF2, 0xF4, 0xF7),
              cell_size=12, head_size=13):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for ci, val in enumerate(data[0]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(val)
        r.font.name = FONT_BODY
        r.font.size = Pt(head_size)
        r.font.bold = True
        r.font.color.rgb = header_color
    for ri in range(1, rows):
        for ci, val in enumerate(data[ri]):
            cell = tbl.cell(ri, ci)
            fill = (alt_fill if ri % 2 == 0 else None) if alt_fill else None
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT_BODY
            r.font.size = Pt(cell_size)
            r.font.color.rgb = GREY_DARK
    return tbl


def add_footer(slide, idx, total):
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
                 "Fund Finance Facility Risk Monitor  |  Synthetic data — illustrative only",
                 size=10, color=GREY_LIGHT)
    add_text_box(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                 f"{idx} / {total}", size=10, color=GREY_LIGHT, align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------
def slide_title(prs):
    s = add_blank_slide(prs)
    # full navy background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # gold accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(1.0), Inches(2.6), Inches(0.15), Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_text_box(s, Inches(1.4), Inches(2.5), Inches(11), Inches(1.2),
                 "Fund Finance Facility", size=44, bold=True, color=WHITE,
                 font=FONT_HEAD)
    add_text_box(s, Inches(1.4), Inches(3.3), Inches(11), Inches(1.0),
                 "Risk Monitor", size=44, bold=True, color=WHITE,
                 font=FONT_HEAD)
    add_text_box(s, Inches(1.4), Inches(4.3), Inches(11), Inches(0.6),
                 "An engineering view of lender-side credit logic",
                 size=22, color=GOLD)
    add_text_box(s, Inches(1.4), Inches(5.1), Inches(11), Inches(0.5),
                 "6 layers  •  4 industry-standard upgrades  •  14 tests",
                 size=16, color=GREY_LIGHT)


def slide_problem(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "The problem",
                    "Why fund finance has no off-the-shelf library")

    bullets = [
        "A bank lends to a private fund — secured by uncalled commitments (subscription line) or by NAV (NAV facility).",
        "Risk question is not 'will the market go up?' but 'how much can we safely lend, where is the binding constraint, and when do we act?'",
        "Stack is fragmented: borrowing base + credit scoring + concentration + liquidity + covenants + stress + triggers. No QuantLib equivalent.",
        "Common interview gap: candidates can talk LTV, but cannot defend why LTV is the wrong covenant to watch first.",
        "This project is a self-contained, modular re-implementation of the underwriting logic — designed to be explained, not productionised.",
    ]
    add_bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.5),
                bullets, size=18)

    # callout box on the right
    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.6), Inches(5.0), Inches(5.2), Inches(1.8))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    callout.line.color.rgb = NAVY
    callout.line.width = Pt(1.5)
    add_text_box(s, Inches(7.85), Inches(5.15), Inches(4.95), Inches(0.4),
                 "Design contribution", size=14, bold=True, color=NAVY)
    add_text_box(s, Inches(7.85), Inches(5.5), Inches(4.95), Inches(1.3),
                 "Six reusable sub-problems, each borrowing a mature open-source / "
                 "industry pattern, recomposed into one engine.",
                 size=14, color=GREY_DARK)


def slide_architecture(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Architecture: six composable layers",
                    "Each layer borrows a mature open-source / industry pattern")

    arch_path = FIGURES / "architecture.png"
    if arch_path.exists():
        # Place image on the left, captions on the right
        pic = add_image(s, arch_path, Inches(0.4), Inches(1.4), height=Inches(5.4))
    else:
        add_text_box(s, Inches(0.7), Inches(3.5), Inches(6), Inches(1),
                     "(architecture.png missing — re-run make_diagram.py)",
                     size=14, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

    layers = [
        ["L5", "Data ingest + validation", "Collateral-mgmt range checks"],
        ["L2", "Credit / EL", "Basel III IRB EL = PD × LGD"],
        ["L1", "Borrowing base + HF", "Aave / Compound Health Factor"],
        ["L5", "Liquidity", "Per-asset haircut by days-to-sell"],
        ["L3", "Stress (3 views)", "Sweep • reverse-per-covenant • multi-factor MC"],
        ["L4", "Covenants + triggers", "Protective-relay with HF-aware cure"],
    ]
    add_table(s, Inches(7.5), Inches(1.4), Inches(5.4), Inches(4.5),
              [["#", "Layer", "Borrowed pattern"]] + layers,
              cell_size=13, head_size=14)


def slide_upgrade1_el(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Upgrade 1 — Basel III IRB expected-loss weighting",
                    "Continuous risk-weighted collateral, not a step function")

    # Left: explanation
    add_text_box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
                 "Formula", size=16, bold=True, color=NAVY)
    add_text_box(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(1.0),
                 "EL = PD × LGD       (Basel III IRB)\n"
                 "weight = max(0, 1 − EL / EL_CAP)",
                 size=18, color=GREY_DARK, font="Consolas")

    add_text_box(s, Inches(0.6), Inches(3.4), Inches(6.0), Inches(0.4),
                 "Why this matters", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(3.0), [
        "v1: binary cliff at PD ≤ 3% — Family Office D excluded entirely.",
        "v2: continuous weight ∈ [0,1] — D contributes 50% (partial), B contributes 99%.",
        "Step changes in eligibility cause step changes in borrowing base.",
        "EE analogue: power-device derating curve, not bang-bang.",
    ], size=15)

    # Right: table of LPs with EL
    add_text_box(s, Inches(7.4), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Per-LP impact (sample fund)", size=16, bold=True, color=NAVY)
    data = [
        ["LP", "Rating", "PD", "LGD", "EL", "Weight"],
        ["Pension Fund A", "A", "0.50%", "40%", "0.20%", "0.98"],
        ["Insurer B", "AA", "0.20%", "35%", "0.07%", "0.99"],
        ["Endowment C", "BBB", "2.00%", "45%", "0.90%", "0.91"],
        ["Family Office D", "NR", "10.00%", "50%", "5.00%", "0.50"],
    ]
    add_table(s, Inches(7.4), Inches(1.95), Inches(5.5), Inches(3.0),
              data, cell_size=12, head_size=13)

    # callout
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.4), Inches(5.2), Inches(5.5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    box.line.color.rgb = NAVY
    add_text_box(s, Inches(7.65), Inches(5.35), Inches(5.0), Inches(0.4),
                 "Subscription line BB", size=13, bold=True, color=NAVY)
    add_text_box(s, Inches(7.65), Inches(5.75), Inches(5.0), Inches(0.9),
                 "v1: 58.5  →  v2: 65.7  (because D now contributes\n"
                 "partially under continuous EL weighting)",
                 size=12, color=GREY_DARK)


def slide_upgrade2_hf(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Upgrade 2 — Health Factor (Aave / Compound DeFi)",
                    "One number, normalised, with a clear liquidation threshold at 1.0")

    add_text_box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
                 "Formula", size=16, bold=True, color=NAVY)
    add_text_box(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(1.0),
                 "HF = (effective_limit × liq_thresh) / drawn\n"
                 "    = (33.75 × 0.85) / 30  =  0.956",
                 size=18, color=GREY_DARK, font="Consolas")

    add_text_box(s, Inches(0.6), Inches(3.4), Inches(6.0), Inches(0.4),
                 "Why this matters", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(3.0), [
        "Utilization has no distance concept — 89% does not tell you how far.",
        "HF is normalised: HF=1.0 is the liquidation threshold, period.",
        "Captures collateral × threshold × debt in one number — sample HF=0.956 → already in margin-call territory.",
        "EE analogue: stability margin (phase/gain margin).",
    ], size=15)

    # Right: HF gauge as bands
    add_text_box(s, Inches(7.4), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Health Factor bands", size=16, bold=True, color=NAVY)
    bands = [
        ("HF ≥ 1.5", "GREEN — safe", GREEN),
        ("1.0 ≤ HF < 1.5", "AMBER — early warning", AMBER),
        ("HF < 1.0", "RED — margin call → default in cure_days", RED),
    ]
    for i, (rng, lbl, color) in enumerate(bands):
        top = Inches(2.0 + i * 0.95)
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(7.4), top, Inches(0.4), Inches(0.85))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        add_text_box(s, Inches(7.95), top + Inches(0.05), Inches(5.0), Inches(0.4),
                     rng, size=14, bold=True, color=GREY_DARK)
        add_text_box(s, Inches(7.95), top + Inches(0.4), Inches(5.0), Inches(0.5),
                     lbl, size=12, color=GREY_DARK)

    # Sample box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.4), Inches(5.0), Inches(5.5), Inches(1.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RED
    box.line.fill.background()
    add_text_box(s, Inches(7.65), Inches(5.15), Inches(5.0), Inches(0.5),
                 "Sample NAV facility today", size=13, bold=True, color=WHITE)
    add_text_box(s, Inches(7.65), Inches(5.6), Inches(5.0), Inches(1.1),
                 "HF = 0.956  →  RED  →  MARGIN CALL\n"
                 "HF-aware cure window = 5 days (vs 10 negotiated)",
                 size=13, color=WHITE)


def slide_upgrade3_mc(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Upgrade 3 — Multi-factor Monte Carlo (MSCI Barra style)",
                    "Sector correlations + VaR + CVaR + per-sector attribution")

    add_text_box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
                 "What changed", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(2.0), [
        "v1: single lognormal × all assets (implicit: sectors 100% correlated).",
        "v2: 10-dim multivariate lognormal with MSCI-Barra-style correlation.",
        "Vectorised with numpy — ~100× speedup vs Python loop.",
        "Adds VaR(95), CVaR(99) and per-sector attribution.",
    ], size=13, line_spacing=1.15)

    add_text_box(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(0.4),
                 "Two numbers, two stories", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(4.6), Inches(6.0), Inches(2.3), [
        "P(LTV breach) drops 22% → 0% — diversification is now real.",
        "But CVaR(99) of LTV = 22.1%, glued to the 25% covenant cap.",
        "Single-shock MC hid both findings under one P(breach) number.",
    ], size=13, line_spacing=1.15)

    # Right: result table
    add_text_box(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(0.4),
                 "Sample run — NAV facility, 6-month horizon", size=16, bold=True, color=NAVY)
    data = [
        ["Metric", "v1 (single shock)", "v2 (multi-factor)"],
        ["P(LTV breach)", "22%", "0.0%"],
        ["VaR(95) of LTV", "(n/a)", "19.4%"],
        ["CVaR(99) of LTV", "(n/a)", "22.1%"],
        ["LTV p50", "16.5%", "15.7%"],
        ["LTV p99", "26.8%", "21.2%"],
        ["Sector attribution", "(n/a)", "Tech 49% / Fin 41% / Ind 32%"],
    ]
    add_table(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(3.6),
              data, cell_size=12, head_size=13)

    # callout
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.2), Inches(5.8), Inches(5.6), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7)
    box.line.color.rgb = NAVY
    add_text_box(s, Inches(7.4), Inches(5.95), Inches(5.2), Inches(0.9),
                 "EE analogue: cross-PSD noise analysis — each sector "
                 "has its own variance, off-diagonal coupling describes how "
                 "they move together.",
                 size=12, color=GREY_DARK)


def slide_upgrade4_reverse(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Upgrade 4 — Reverse stress per covenant (Basel SR 11-7)",
                    "Distance-to-breach ranked across every covenant — stability margin per loop")

    add_text_box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
                 "Algorithm", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(2.0), [
        "Bisection on NAV haircut for every covenant.",
        "O(covenants × log(1/tol)) — fast even at scale.",
        "Returns the smallest haircut that flips each covenant to RED.",
    ], size=14)

    add_text_box(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(0.4),
                 "Headline finding (sample fund)", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(4.45), Inches(6.0), Inches(2.5), [
        "Health Factor binds today (0% haircut).",
        "BB and Utilization both at 11.2% NAV drop.",
        "Max LTV at 36.9% — 25 percentage points further away.",
        "Asset / sector concentration are invariant under uniform NAV drop (need single-sector shock).",
    ], size=14)

    # Right: image
    rs_path = FIGURES / "reverse_stress.png"
    if rs_path.exists():
        add_image(s, rs_path, Inches(6.8), Inches(1.5), width=Inches(6.4))
    else:
        add_text_box(s, Inches(6.8), Inches(3.5), Inches(6.0), Inches(0.6),
                     "(reverse_stress.png missing — run run.py first)",
                     size=14, color=GREY_LIGHT)


def slide_sample_run(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Sample run — NAV facility on Project Alpha Fund",
                    "Four eligible LPs, five sector-diverse assets, intentionally tight")

    # Left numerical summary
    add_text_box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.4),
                 "Facility state today", size=16, bold=True, color=NAVY)
    data = [
        ["Metric", "Value", "Status"],
        ["Borrowing base", "33.75", "after liquidity haircut"],
        ["Drawn", "30.00", ""],
        ["Availability", "3.75", "very tight"],
        ["Utilization", "88.9%", "AMBER"],
        ["LTV", "15.8%", "well under 25% cap"],
        ["Health Factor", "0.956", "RED — margin call"],
        ["Liquidity coverage", "5.17x", "GREEN"],
    ]
    add_table(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(4.0),
              data, cell_size=12, head_size=13)

    # Right reverse stress ranking
    add_text_box(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
                 "Distance-to-breach ranking", size=16, bold=True, color=NAVY)
    rs = [
        ["Rank", "Covenant", "Breaks at NAV drop"],
        ["1", "Health Factor", "0.0%  ← BINDING TODAY"],
        ["2", "BB headroom", "11.2%"],
        ["3", "Utilization", "11.2%"],
        ["4", "Max LTV", "36.9%"],
        ["5", "Liquidity coverage", "80.7%"],
        ["—", "Asset / sector / #assets", "never under uniform haircut"],
    ]
    add_table(s, Inches(7.0), Inches(1.95), Inches(5.8), Inches(3.5),
              rs, cell_size=12, head_size=13)

    # Bottom takeaway band
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, Inches(6.2), SLIDE_W, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text_box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                 "Binding constraint is the Health Factor — not the LTV.  "
                 "Without reverse stress per covenant you would miss this.",
                 size=15, bold=True, color=WHITE)


def slide_three_views(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Three stress views, one engine",
                    "Step input, ramp input, noise input — through the same plant")

    # left: stress_chart.png
    sc_path = FIGURES / "stress_chart.png"
    if sc_path.exists():
        add_image(s, sc_path, Inches(0.4), Inches(1.4), width=Inches(6.3))
    add_text_box(s, Inches(0.4), Inches(5.9), Inches(6.3), Inches(0.4),
                 "Deterministic sweep — LTV and availability vs NAV haircut",
                 size=12, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

    # right: reverse stress + MC summary
    rs_path = FIGURES / "reverse_stress.png"
    if rs_path.exists():
        add_image(s, rs_path, Inches(6.9), Inches(1.4), width=Inches(6.1))
    add_text_box(s, Inches(6.9), Inches(5.9), Inches(6.1), Inches(0.4),
                 "Reverse stress — distance-to-breach per covenant",
                 size=12, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

    # bottom band
    add_text_box(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.6),
                 "Multi-factor MC adds the third view: probability of breach, "
                 "VaR(95), CVaR(99) and which sector is driving the tail. "
                 "Three independent views; never trust one alone.",
                 size=14, color=GREY_DARK, align=PP_ALIGN.CENTER)


def slide_isomorphism(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Engineering ↔ fund finance isomorphism",
                    "Why an EE perspective is structural, not decorative")

    data = [
        ["Fund-finance concept", "EE concept", "Same problem solved"],
        ["Advance rate", "Derating factor", "Margin against operating point"],
        ["EL weight (PD × LGD)", "Continuous gain function", "Avoid step changes in main loop"],
        ["Liquidity haircut", "Bandwidth-limited signal", "Discount slow / noisy inputs"],
        ["Borrowing base", "Plant capacity estimate", "Safe operating envelope"],
        ["Health Factor", "Stability margin", "Normalised distance to instability"],
        ["Covenant alarm (AMBER)", "Pre-trip alarm", "Detect before fault"],
        ["Covenant breach (RED)", "Trip threshold", "Open the breaker"],
        ["Cure period", "Time-graded trip delay", "Avoid nuisance trips"],
        ["HF-aware cure (shorter when worse)", "Inverse-time overcurrent relay", "Faster trip closer to fault"],
        ["Reverse stress per covenant", "Per-loop stability margin", "Find the weakest loop"],
        ["Multi-factor MC attribution", "Cross-PSD noise analysis", "Decompose disturbance source"],
        ["Margin call → default", "Two-stage breaker", "Graduated protection"],
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4),
              data, cell_size=11, head_size=12)


def slide_pitch(prs):
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    add_text_box(s, Inches(0.7), Inches(0.6), Inches(12), Inches(0.6),
                 "30-second pitch", size=22, bold=True, color=GOLD)

    pitch = (
        '"I built a fund-finance facility risk monitor. Because fund finance '
        'is a fusion problem with no ready library, I decomposed it into six '
        'layers — data validation, a Basel-IRB-style expected-loss weighting '
        'on subscription collateral, a borrowing base with an Aave-style Health '
        'Factor, reverse-stress-per-covenant ranking, and a multi-factor Monte '
        'Carlo with sector correlations giving VaR and CVaR.\n\n'
        'My EE background means I treat each piece as one part of a feedback '
        'control system: covenants are protective relays with an alarm band '
        'and a timed trip, the advance rate is a derating curve not a cliff, '
        'the Health Factor is a stability margin, and reverse stress is per-loop '
        'stability ranking.\n\n'
        'One result on the synthetic fund: a single-shock Monte Carlo says 22% '
        'breach probability; the multi-factor Monte Carlo says effectively zero '
        '— but CVaR(99) on LTV is still 22%, right at the covenant. Two views, '
        'two stories, both true. A lender’s view: protect principal, rank '
        'constraints, don’t forecast returns."'
    )

    box = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = pitch
    r.font.name = FONT_BODY
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE


def slide_extend(prs):
    s = add_blank_slide(prs)
    add_header_band(s, "Ways to extend",
                    "Roadmap items, each with a known open-source / industry anchor")

    items = [
        ("Trained PD model", "Replace rating-table lookup with logistic regression or XGBoost trained on borrower features — interface already exposed at credit_scoring.estimate_pd."),
        ("Multi-period MC", "Daily random-walk over the horizon; demonstrates the HF-aware cure window numerically; matches CCAR base / adverse / severe scenario libraries."),
        ("LP-default scenarios", "drop_lp() helper plus reverse-stress engine → 'what NAV drop breaks each covenant if LP X defaults?'"),
        ("Per-sector reverse stress", "Single-sector shock instead of uniform haircut → asset / sector concentration covenants become reachable."),
        ("Streamlit / FastAPI dashboard", "Same engine, web UI, live re-calculation per facility, ranked daily."),
        ("Historical calibration", "Estimate sector vols and correlations from real PE-fund NAV history; replace MSCI-Barra approximations."),
    ]
    top = Inches(1.4)
    for title, body in items:
        add_text_box(s, Inches(0.5), top, Inches(12.3), Inches(0.35),
                     title, size=14, bold=True, color=NAVY)
        add_text_box(s, Inches(0.5), top + Inches(0.35), Inches(12.3), Inches(0.45),
                     body, size=11, color=GREY_DARK)
        top += Inches(0.85)


# ----------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_architecture,
        slide_upgrade1_el,
        slide_upgrade2_hf,
        slide_upgrade3_mc,
        slide_upgrade4_reverse,
        slide_sample_run,
        slide_three_views,
        slide_isomorphism,
        slide_pitch,
        slide_extend,
    ]
    total = len(builders)
    for idx, build in enumerate(builders, start=1):
        build(prs)
        if idx != 1 and idx != 11:    # skip footer on full-bleed slides
            add_footer(prs.slides[idx - 1], idx, total)

    SLIDES.mkdir(exist_ok=True)
    out_path = SLIDES / "FFRM-Deck.pptx"
    prs.save(str(out_path))
    print(f"Saved {out_path} ({total} slides)")


if __name__ == "__main__":
    main()
